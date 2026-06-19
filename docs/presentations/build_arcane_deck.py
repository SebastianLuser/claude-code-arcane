#!/usr/bin/env python3
"""Genera la presentacion ejecutiva de Claude Code Arcane (8 slides, 16:9)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- Tema (paleta Arcane: indigo profundo + cian) ----------
BG       = RGBColor(0x1E, 0x1B, 0x4B)  # indigo-950 fondo
PANEL    = RGBColor(0x2A, 0x25, 0x60)  # panel
PANEL2   = RGBColor(0x31, 0x2E, 0x81)  # panel mas claro
VIOLET   = RGBColor(0x5B, 0x4F, 0xE8)  # acento violeta
CYAN     = RGBColor(0x22, 0xD3, 0xEE)  # acento cian
WHITE    = RGBColor(0xF8, 0xFA, 0xFC)
LIGHT    = RGBColor(0xE5, 0xE7, 0xEB)
MUTED    = RGBColor(0x9C, 0xA3, 0xC4)
GREEN    = RGBColor(0x34, 0xD3, 0x99)
AMBER    = RGBColor(0xFB, 0xBF, 0x24)
FONT     = "Calibri"
FONT_H   = "Calibri"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def rect(s, x, y, w, h, color, line=None, radius=False, line_w=1.0):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sp = s.shapes.add_shape(shp, Inches(x), Inches(y), Inches(w), Inches(h))
    if color is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=4, line_spacing=1.0):
    """runs: lista de parrafos; cada parrafo es lista de (texto, size, color, bold)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after); p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, sz, col, bold) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col
            r.font.bold = bold; r.font.name = FONT
    return tb


def badge(s, n):
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(12.55), Inches(6.85), Inches(0.42), Inches(0.42))
    c.fill.solid(); c.fill.fore_color.rgb = VIOLET; c.line.fill.background(); c.shadow.inherit = False
    tf = c.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n)
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def footer(s):
    txt(s, 0.55, 7.0, 5, 0.35, [[("Claude Code Arcane", 9, MUTED, False)]])


def kicker(s, text, color=CYAN):
    rect(s, 0.55, 0.55, 0.12, 0.55, color)
    txt(s, 0.8, 0.5, 11, 0.5, [[(text, 13, color, True)]])


def title(s, text, y=0.95, size=34):
    txt(s, 0.78, y, 11.8, 1.0, [[(text, size, WHITE, True)]])


# ================= SLIDE 1 — Portada =================
s = slide()
rect(s, 0, 0, 13.333, 7.5, None)
# barra de acento lateral
rect(s, 0, 0, 0.25, 7.5, VIOLET)
rect(s, 0.25, 0, 0.08, 7.5, CYAN)
txt(s, 1.0, 1.7, 11.5, 0.6, [[("CLAUDE CODE", 20, CYAN, True)]])
txt(s, 1.0, 2.15, 11.5, 1.4, [[("Arcane", 88, WHITE, True)]])
txt(s, 1.05, 3.65, 11, 0.6,
    [[("Un gestor de paquetes de configuracion para Claude Code", 22, LIGHT, False)]])
txt(s, 1.05, 4.25, 11, 0.5,
    [[("Instala solo las capacidades que tu proyecto necesita — con un comando.", 15, MUTED, False)]])

# dashboard de metricas
metrics = [("322", "Skills"), ("86", "Agentes"), ("29", "Profiles"), ("19", "Rules")]
mx, mw, gap = 1.0, 2.7, 0.25
for i, (num, lbl) in enumerate(metrics):
    x = mx + i * (mw + gap)
    rect(s, x, 5.3, mw, 1.35, PANEL, radius=True)
    rect(s, x, 5.3, 0.1, 1.35, CYAN if i % 2 else VIOLET, radius=False)
    txt(s, x, 5.45, mw, 0.8, [[(num, 44, WHITE, True)]], align=PP_ALIGN.CENTER)
    txt(s, x, 6.2, mw, 0.4, [[(lbl.upper(), 13, CYAN, True)]], align=PP_ALIGN.CENTER)

# ================= SLIDE 2 — El problema =================
s = slide()
kicker(s, "POR QUE EXISTE")
title(s, "El problema que resuelve")
problems = [
    ("Configuracion manual", "Copiar .claude/ a mano en cada proyecto: repetitivo y diverge con el tiempo."),
    ("Todo-o-nada", "Un mega-repo cargaria las 322 skills en CADA proyecto: mas tokens, mas ruido, mas latencia."),
    ("Sin estandarizacion", "Cada dev configura su Claude distinto. No hay fuente de verdad compartida."),
    ("Sin actualizacion", "Una mejora en una skill nunca llega a los proyectos que ya la copiaron."),
]
y = 2.05
for i, (h, d) in enumerate(problems):
    rect(s, 0.78, y, 11.7, 1.0, PANEL, radius=True)
    rect(s, 0.78, y, 0.1, 1.0, VIOLET)
    txt(s, 1.05, y + 0.13, 3.4, 0.8, [[(h, 17, CYAN, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 4.5, y + 0.13, 7.8, 0.8, [[(d, 13.5, LIGHT, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += 1.15
txt(s, 0.78, 6.75, 11.8, 0.5,
    [[("Falta una forma modular, versionada y compartible de configurar Claude Code.", 15, WHITE, True)]])
badge(s, 2); footer(s)

# ================= SLIDE 3 — Idea central =================
s = slide()
kicker(s, "LA IDEA CENTRAL")
title(s, "npm, pero para las capacidades de Claude Code")
txt(s, 0.78, 2.0, 11.7, 1.0,
    [[("En vez de instalar librerias de codigo, instalas ", 17, LIGHT, False),
      ("capacidades", 17, CYAN, True),
      (": hacer commits, disenar APIs, correr un sprint, auditar seguridad — adaptadas a tu stack.", 17, LIGHT, False)]],
    line_spacing=1.15)

# comando hero estilo terminal
rect(s, 0.78, 3.25, 11.7, 1.5, RGBColor(0x14, 0x11, 0x33), radius=True, line=VIOLET, line_w=1.5)
txt(s, 1.1, 3.45, 11, 0.4, [[("terminal", 11, MUTED, False)]])
txt(s, 1.1, 3.85, 11, 0.7,
    [[("$ ", 26, CYAN, True), ("npx arcane install ", 26, WHITE, True),
      ("backend-ts", 26, GREEN, True), ("+agile", 26, AMBER, True), ("+testing", 26, AMBER, True)]])

cols = [("backend-ts", "base", GREEN), ("+agile", "addon", AMBER), ("+testing", "addon", AMBER)]
cx = 0.78
for name, kind, col in cols:
    rect(s, cx, 5.05, 3.0, 1.1, PANEL, radius=True)
    txt(s, cx, 5.2, 3.0, 0.5, [[(name, 18, col, True)]], align=PP_ALIGN.CENTER)
    txt(s, cx, 5.7, 3.0, 0.4, [[(kind.upper(), 12, MUTED, True)]], align=PP_ALIGN.CENTER)
    cx += 3.25
txt(s, 0.78, 6.45, 11.7, 0.5,
    [[("base + addon + addon  ->  solo lo que ESE proyecto necesita.", 15, WHITE, True)]])
badge(s, 3); footer(s)

# ================= SLIDE 4 — Como funciona =================
s = slide()
kicker(s, "ARQUITECTURA")
title(s, "Como funciona")
steps = [
    ("1", "Resuelve el contenido", "GitHub -> cache local -> bundled (funciona offline)"),
    ("2", "Carga core.yaml siempre", "21 skills base + 15 hooks + permisos de seguridad"),
    ("3", "Mergea los profiles", "Combina los profiles pedidos con +"),
    ("4", "Deduplica", "Un skill repetido se copia una sola vez"),
    ("5", "Genera .claude/", "settings.json + manifest + skills, agents, rules, hooks"),
]
y = 2.05
for num, h, d in steps:
    rect(s, 0.78, y, 11.7, 0.84, PANEL, radius=True)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.98), Inches(y + 0.17), Inches(0.5), Inches(0.5))
    c.fill.solid(); c.fill.fore_color.rgb = VIOLET; c.line.fill.background(); c.shadow.inherit = False
    ctf = c.text_frame; ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run(); cr.text = num; cr.font.size = Pt(20); cr.font.bold = True
    cr.font.color.rgb = WHITE; cr.font.name = FONT
    txt(s, 1.75, y + 0.05, 4.0, 0.74, [[(h, 16, CYAN, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 5.9, y + 0.05, 6.4, 0.74, [[(d, 13, LIGHT, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.95
badge(s, 4); footer(s)

# ================= SLIDE 5 — 4 piezas =================
s = slide()
kicker(s, "SEPARACION DE RESPONSABILIDADES")
title(s, "Cuatro piezas, cuatro responsabilidades")
cards = [
    ("Skills", "Acciones que Claude ejecuta", "/commit  ·  /code-review", CYAN),
    ("Agents", "Personas que Claude adopta", "backend-architect  ·  game-designer", VIOLET),
    ("Hooks", "Automatizacion invisible (eventos)", "valida commits  ·  escanea secrets", GREEN),
    ("Rules", "Estandares que Claude obedece", "\"Clean architecture en backend\"", AMBER),
]
positions = [(0.78, 2.05), (6.93, 2.05), (0.78, 4.25), (6.93, 4.25)]
cw, ch = 5.6, 2.0
for (name, what, ex, col), (x, y) in zip(cards, positions):
    rect(s, x, y, cw, ch, PANEL, radius=True)
    rect(s, x, y, cw, 0.12, col)
    txt(s, x + 0.3, y + 0.25, cw - 0.6, 0.5, [[(name, 24, col, True)]])
    txt(s, x + 0.3, y + 0.85, cw - 0.6, 0.5, [[(what, 14.5, LIGHT, False)]])
    txt(s, x + 0.3, y + 1.4, cw - 0.6, 0.5, [[(ex, 13, MUTED, True)]])
txt(s, 0.78, 6.55, 11.7, 0.5,
    [[("Todo en Markdown + YAML: legible, diffeable en git, editable por cualquiera.", 14, WHITE, True)]])
badge(s, 5); footer(s)

# ================= SLIDE 6 — Por que gana =================
s = slide()
kicker(s, "DIFERENCIADORES")
title(s, "Por que Arcane gana")
# encabezados
rect(s, 0.78, 2.0, 5.7, 0.55, PANEL2, radius=True)
rect(s, 6.63, 2.0, 5.85, 0.55, PANEL2, radius=True)
txt(s, 1.0, 2.05, 5.3, 0.45, [[("ALTERNATIVA", 13, MUTED, True)]], anchor=MSO_ANCHOR.MIDDLE)
txt(s, 6.85, 2.05, 5.4, 0.45, [[("COMO GANA ARCANE", 13, CYAN, True)]], anchor=MSO_ANCHOR.MIDDLE)
rows = [
    ("Config a mano", "Un comando, versionado, reproducible"),
    ("Mega-repo monolitico", "Selectivo por profile + dedup (menos tokens)"),
    ("Copiar config entre proyectos", "Smart update con three-way merge (no pisa tus cambios)"),
    ("Claude \"vanilla\"", "Hooks + rules + 86 agentes especializados"),
]
y = 2.65
for a, b in rows:
    rect(s, 0.78, y, 5.7, 0.7, PANEL, radius=True)
    rect(s, 6.63, y, 5.85, 0.7, PANEL, radius=True)
    txt(s, 1.0, y + 0.05, 5.3, 0.6, [[(a, 14, LIGHT, False)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 6.85, y + 0.05, 5.4, 0.6, [[(b, 13.5, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.8
rect(s, 0.78, y + 0.05, 11.7, 0.7, VIOLET, radius=True)
txt(s, 0.78, y + 0.1, 11.7, 0.6,
    [[("Trata la configuracion de IA como software de primera clase: versionado, modular, testeable y distribuido.",
       14, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
badge(s, 6); footer(s)

# ================= SLIDE 7 — Pros y contras =================
s = slide()
kicker(s, "EVALUACION HONESTA")
title(s, "Trade-offs")
# Fortalezas
rect(s, 0.78, 2.0, 5.7, 4.6, PANEL, radius=True)
rect(s, 0.78, 2.0, 5.7, 0.6, RGBColor(0x14, 0x3D, 0x2E), radius=True)
txt(s, 1.05, 2.05, 5.3, 0.5, [[("FORTALEZAS", 16, GREEN, True)]], anchor=MSO_ANCHOR.MIDDLE)
pros = [
    "Eficiencia de tokens (solo lo relevante)",
    "Un comando, cero friccion (npx)",
    "Reproducible (manifest = snapshot exacto)",
    "Estandarizacion de equipo",
    "Auto-actualizable sin pisar tus cambios",
    "Testeable (framework skills-selftest)",
]
y = 2.85
for p in pros:
    txt(s, 1.05, y, 5.2, 0.5, [[("+  ", 14, GREEN, True), (p, 13.5, LIGHT, False)]])
    y += 0.6
# Costos
rect(s, 6.93, 2.0, 5.55, 4.6, PANEL, radius=True)
rect(s, 6.93, 2.0, 5.55, 0.6, RGBColor(0x4A, 0x39, 0x12), radius=True)
txt(s, 7.2, 2.05, 5.1, 0.5, [[("COSTOS / TRADE-OFFS", 16, AMBER, True)]], anchor=MSO_ANCHOR.MIDDLE)
cons = [
    "Mantener 322 skills con calidad consistente",
    "Skills son prompts -> no deterministas",
    "Dependencia de la API de Claude Code",
    "Curva de aprendizaje para combinar profiles",
]
y = 2.95
for c in cons:
    txt(s, 7.2, y, 5.0, 0.6, [[("!  ", 14, AMBER, True), (c, 13.5, LIGHT, False)]])
    y += 0.75
badge(s, 7); footer(s)

# ================= SLIDE 8 — Cierre =================
s = slide()
rect(s, 0, 0, 0.25, 7.5, VIOLET)
rect(s, 0.25, 0, 0.08, 7.5, CYAN)
kicker(s, "PRODUCTION-GRADE")
title(s, "No es un script: es infraestructura", size=32)
feats = [
    ("Distribucion hibrida", "GitHub -> cache -> bundled. Funciona incluso offline."),
    ("Smart update", "Three-way merge: actualiza sin pisar tus customizaciones."),
    ("Worktree-aware", "Comparte hooks/docs, copia skills por rama."),
    ("Auto-release", "semantic-release: conventional commits -> publish a npm."),
]
positions = [(0.78, 2.15), (6.93, 2.15), (0.78, 3.75), (6.93, 3.75)]
for (h, d), (x, y) in zip(feats, positions):
    rect(s, x, y, 5.6, 1.4, PANEL, radius=True)
    rect(s, x, y, 0.1, 1.4, CYAN)
    txt(s, x + 0.3, y + 0.18, 5.1, 0.5, [[(h, 17, CYAN, True)]])
    txt(s, x + 0.3, y + 0.68, 5.1, 0.6, [[(d, 13, LIGHT, False)]])

rect(s, 0.78, 5.55, 11.7, 1.3, VIOLET, radius=True)
txt(s, 0.78, 5.7, 11.7, 0.55,
    [[("Configuracion de IA, modular y versionada — un comando para empezar.", 19, WHITE, True)]],
    align=PP_ALIGN.CENTER)
txt(s, 0.78, 6.25, 11.7, 0.5,
    [[("npx arcane install <profile>", 22, WHITE, True)]], align=PP_ALIGN.CENTER)
badge(s, 8); footer(s)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "arcane-overview.pptx")
prs.save(out)
print("OK ->", out)
print("Slides:", len(prs.slides._sldIdLst))

#!/usr/bin/env python3
"""Genera la presentacion DETALLADA de Claude Code Arcane (12 slides, 16:9)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------- Tema (paleta Arcane: indigo profundo + cian) ----------
BG     = RGBColor(0x1E, 0x1B, 0x4B)
PANEL  = RGBColor(0x2A, 0x25, 0x60)
PANEL2 = RGBColor(0x31, 0x2E, 0x81)
VIOLET = RGBColor(0x5B, 0x4F, 0xE8)
CYAN   = RGBColor(0x22, 0xD3, 0xEE)
WHITE  = RGBColor(0xF8, 0xFA, 0xFC)
LIGHT  = RGBColor(0xE5, 0xE7, 0xEB)
MUTED  = RGBColor(0x9C, 0xA3, 0xC4)
GREEN  = RGBColor(0x34, 0xD3, 0x99)
AMBER  = RGBColor(0xFB, 0xBF, 0x24)
PINK   = RGBColor(0xF4, 0x72, 0xB6)
CODEBG = RGBColor(0x14, 0x11, 0x33)
FONT   = "Calibri"
MONO   = "Consolas"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background(); r.shadow.inherit = False
    return s


def rect(s, x, y, w, h, color, line=None, radius=False, line_w=1.0):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sp = s.shapes.add_shape(shp, Inches(x), Inches(y), Inches(w), Inches(h))
    if color is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=4, line_spacing=1.0, font=FONT):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after); p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, sz, col, bold) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col
            r.font.bold = bold; r.font.name = font
    return tb


def badge(s, n):
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(12.55), Inches(6.85), Inches(0.42), Inches(0.42))
    c.fill.solid(); c.fill.fore_color.rgb = VIOLET; c.line.fill.background(); c.shadow.inherit = False
    tf = c.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n)
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def footer(s):
    txt(s, 0.55, 7.02, 6, 0.35, [[("Claude Code Arcane", 9, MUTED, False)]])


def kicker(s, text, color=CYAN):
    rect(s, 0.55, 0.55, 0.12, 0.5, color)
    txt(s, 0.8, 0.5, 11.5, 0.45, [[(text, 12.5, color, True)]])


def title(s, text, y=0.92, size=32):
    txt(s, 0.78, y, 11.9, 0.9, [[(text, size, WHITE, True)]])


def chrome(s, n, kick, ttl, color=CYAN, tsize=32):
    kicker(s, kick, color)
    title(s, ttl, size=tsize)
    badge(s, n); footer(s)


# ================= SLIDE 1 — Portada =================
s = slide()
rect(s, 0, 0, 0.25, 7.5, VIOLET)
rect(s, 0.25, 0, 0.08, 7.5, CYAN)
txt(s, 1.0, 1.55, 11.5, 0.6, [[("CLAUDE CODE", 20, CYAN, True)]])
txt(s, 1.0, 2.0, 11.5, 1.4, [[("Arcane", 84, WHITE, True)]])
txt(s, 1.05, 3.42, 11, 0.6, [[("Un gestor de paquetes de configuracion para Claude Code", 21, LIGHT, False)]])
txt(s, 1.05, 4.0, 11, 0.5, [[("Skills, agentes, hooks y reglas — instalados selectivamente por stack.", 14, MUTED, False)]])
metrics = [("322", "Skills"), ("86", "Agentes"), ("29", "Profiles"), ("19", "Rules"), ("16", "Hooks")]
mx, mw, gap = 1.0, 2.12, 0.18
for i, (num, lbl) in enumerate(metrics):
    x = mx + i * (mw + gap)
    rect(s, x, 4.95, mw, 1.4, PANEL, radius=True)
    rect(s, x, 4.95, 0.09, 1.4, CYAN if i % 2 else VIOLET)
    txt(s, x, 5.12, mw, 0.7, [[(num, 40, WHITE, True)]], align=PP_ALIGN.CENTER)
    txt(s, x, 5.88, mw, 0.4, [[(lbl.upper(), 12, CYAN, True)]], align=PP_ALIGN.CENTER)
txt(s, 1.0, 6.7, 11, 0.4, [[("Presentacion tecnica + de negocio  ·  Educabot", 12, MUTED, False)]])

# ================= SLIDE 2 — Indice / recorrido =================
s = slide()
chrome(s, 2, "RECORRIDO", "Que vamos a ver")
items = [
    ("01", "El problema", "Por que la config de Claude no escala"),
    ("02", "La idea", "npm, pero para capacidades de IA"),
    ("03", "Arquitectura", "Como instala, mergea y deduplica"),
    ("04", "Las 4 piezas", "Skills, Agents, Hooks, Rules"),
    ("05", "Anatomia de una skill", "Frontmatter + markdown"),
    ("06", "Profiles", "El sistema de combinacion LEGO"),
    ("07", "Mecanismos avanzados", "Distribucion, update, worktrees"),
    ("08", "Calidad", "Testing automatizado de skills"),
    ("09", "Por que gana", "vs. las alternativas"),
    ("10", "Trade-offs y cierre", "Honesto + casos reales"),
]
positions = [(0.78, 2.0 + (i % 5) * 0.95, 0 if i < 5 else 1) for i in range(10)]
for (num, h, d), i in zip(items, range(10)):
    col = 0 if i < 5 else 1
    x = 0.78 + col * 6.15
    y = 2.0 + (i % 5) * 0.95
    rect(s, x, y, 5.85, 0.82, PANEL, radius=True)
    txt(s, x + 0.25, y, 0.9, 0.82, [[(num, 22, VIOLET, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + 1.15, y + 0.1, 4.5, 0.4, [[(h, 14.5, CYAN, True)]])
    txt(s, x + 1.15, y + 0.45, 4.6, 0.35, [[(d, 11, MUTED, False)]])

# ================= SLIDE 3 — El problema =================
s = slide()
chrome(s, 3, "POR QUE EXISTE", "El problema que resuelve")
problems = [
    ("Configuracion manual", "Copiar .claude/ a mano en cada proyecto: repetitivo y diverge con el tiempo."),
    ("Todo-o-nada", "Un mega-repo cargaria las 322 skills en CADA proyecto: mas tokens, mas ruido, mas latencia."),
    ("Sin estandarizacion", "Cada dev configura su Claude distinto. No hay fuente de verdad compartida."),
    ("Sin actualizacion", "Una mejora en una skill nunca llega a los proyectos que ya la copiaron."),
]
y = 2.0
for h, d in problems:
    rect(s, 0.78, y, 11.7, 0.95, PANEL, radius=True)
    rect(s, 0.78, y, 0.1, 0.95, VIOLET)
    txt(s, 1.05, y + 0.1, 3.4, 0.75, [[(h, 16, CYAN, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 4.5, y + 0.1, 7.8, 0.75, [[(d, 13, LIGHT, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += 1.08
txt(s, 0.78, 6.55, 11.8, 0.5,
    [[("Falta una forma modular, versionada y compartible de configurar Claude Code.", 15, WHITE, True)]])

# ================= SLIDE 4 — La idea central =================
s = slide()
chrome(s, 4, "LA IDEA CENTRAL", "npm, pero para las capacidades de Claude Code")
txt(s, 0.78, 1.95, 11.7, 0.9,
    [[("En vez de instalar librerias de codigo, instalas ", 16, LIGHT, False),
      ("capacidades", 16, CYAN, True),
      (": hacer commits, disenar APIs, correr un sprint, auditar seguridad — adaptadas a tu stack.", 16, LIGHT, False)]],
    line_spacing=1.12)
rect(s, 0.78, 3.05, 11.7, 1.45, CODEBG, radius=True, line=VIOLET, line_w=1.5)
txt(s, 1.1, 3.22, 11, 0.35, [[("terminal", 11, MUTED, False)]], font=MONO)
txt(s, 1.1, 3.62, 11, 0.7,
    [[("$ ", 24, CYAN, True), ("npx arcane install ", 24, WHITE, True),
      ("backend-ts", 24, GREEN, True), ("+agile", 24, AMBER, True), ("+testing", 24, AMBER, True)]], font=MONO)
cols = [("backend-ts", "base", GREEN), ("+agile", "addon", AMBER), ("+testing", "addon", AMBER)]
cx = 0.78
for name, kind, col in cols:
    rect(s, cx, 4.75, 3.0, 1.05, PANEL, radius=True)
    txt(s, cx, 4.9, 3.0, 0.5, [[(name, 17, col, True)]], align=PP_ALIGN.CENTER)
    txt(s, cx, 5.4, 3.0, 0.35, [[(kind.upper(), 11, MUTED, True)]], align=PP_ALIGN.CENTER)
    cx += 3.25
txt(s, 0.78, 6.05, 11.7, 0.5, [[("base + addon + addon  ->  solo lo que ESE proyecto necesita.", 15, WHITE, True)]])

# ================= SLIDE 5 — Como funciona =================
s = slide()
chrome(s, 5, "ARQUITECTURA", "Como funciona la instalacion")
steps = [
    ("1", "Resuelve el contenido", "GitHub -> cache local -> bundled (funciona offline)"),
    ("2", "Carga core.yaml siempre", "21 skills base + 15 hooks + permisos de seguridad"),
    ("3", "Mergea los profiles", "Combina los profiles pedidos con el separador +"),
    ("4", "Deduplica", "Un skill repetido en varios profiles se copia una sola vez"),
    ("5", "Genera .claude/", "settings.json + manifest + skills, agents, rules, hooks, docs"),
]
y = 1.95
for num, h, d in steps:
    rect(s, 0.78, y, 11.7, 0.82, PANEL, radius=True)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.98), Inches(y + 0.16), Inches(0.5), Inches(0.5))
    c.fill.solid(); c.fill.fore_color.rgb = VIOLET; c.line.fill.background(); c.shadow.inherit = False
    ctf = c.text_frame; ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run(); cr.text = num; cr.font.size = Pt(20); cr.font.bold = True
    cr.font.color.rgb = WHITE; cr.font.name = FONT
    txt(s, 1.75, y + 0.04, 4.0, 0.74, [[(h, 15.5, CYAN, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 5.9, y + 0.04, 6.4, 0.74, [[(d, 12.5, LIGHT, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.93
txt(s, 0.78, 6.6, 11.7, 0.4,
    [[("Resultado: el .claude/ del proyecto contiene exactamente lo necesario — nada mas.", 13, MUTED, True)]])

# ================= SLIDE 6 — 4 piezas =================
s = slide()
chrome(s, 6, "SEPARACION DE RESPONSABILIDADES", "Cuatro piezas, cuatro responsabilidades")
cards = [
    ("Skills", "Acciones que Claude ejecuta", "/commit  ·  /code-review  ·  /sprint-plan", CYAN),
    ("Agents", "Personas que Claude adopta", "backend-architect  ·  game-designer", VIOLET),
    ("Hooks", "Automatizacion invisible (eventos)", "valida commits  ·  escanea secrets", GREEN),
    ("Rules", "Estandares que Claude obedece", "\"Clean architecture en backend\"", AMBER),
]
positions = [(0.78, 1.95), (6.93, 1.95), (0.78, 4.15), (6.93, 4.15)]
cw, ch = 5.6, 2.0
for (name, what, ex, col), (x, y) in zip(cards, positions):
    rect(s, x, y, cw, ch, PANEL, radius=True)
    rect(s, x, y, cw, 0.12, col)
    txt(s, x + 0.3, y + 0.22, cw - 0.6, 0.5, [[(name, 23, col, True)]])
    txt(s, x + 0.3, y + 0.82, cw - 0.6, 0.5, [[(what, 14, LIGHT, False)]])
    txt(s, x + 0.3, y + 1.35, cw - 0.6, 0.5, [[(ex, 12, MUTED, True)]])
txt(s, 0.78, 6.45, 11.7, 0.5,
    [[("Todo en Markdown + YAML: legible, diffeable en git, editable por cualquiera.", 14, WHITE, True)]])

# ================= SLIDE 7 — Anatomia de una skill =================
s = slide()
chrome(s, 7, "DEEP DIVE", "Anatomia de una skill")
# bloque de codigo
rect(s, 0.78, 1.95, 6.7, 4.55, CODEBG, radius=True, line=VIOLET, line_w=1.2)
txt(s, 1.0, 2.1, 6.3, 0.35, [[("skills/commit/SKILL.md", 11, MUTED, False)]], font=MONO)
code = [
    [("---", 12.5, MUTED, False)],
    [("name:", 12.5, PINK, True), (" commit", 12.5, LIGHT, False)],
    [("description:", 12.5, PINK, True), (" \"Crea un commit", 12.5, GREEN, False)],
    [("  siguiendo conventional commits...\"", 12.5, GREEN, False)],
    [("category:", 12.5, PINK, True), (" workflow", 12.5, LIGHT, False)],
    [("user-invocable:", 12.5, PINK, True), (" true", 12.5, AMBER, False)],
    [("allowed-tools:", 12.5, PINK, True), (" Read, Bash", 12.5, LIGHT, False)],
    [("model:", 12.5, PINK, True), (" haiku", 12.5, LIGHT, False), ("    # opcional", 11, MUTED, False)],
    [("---", 12.5, MUTED, False)],
    [("", 6, MUTED, False)],
    [("Review staged & unstaged changes", 12.5, CYAN, False)],
    [("with git diff, suggest type/scope,", 12.5, CYAN, False)],
    [("write a concise message...", 12.5, CYAN, False)],
]
txt(s, 1.0, 2.5, 6.3, 3.9, code, space_after=2, font=MONO)
# explicacion lado derecho
rect(s, 7.75, 1.95, 4.72, 4.55, PANEL, radius=True)
txt(s, 8.0, 2.15, 4.3, 0.4, [[("Dos partes:", 16, WHITE, True)]])
fields = [
    ("Frontmatter (YAML)", "Metadatos: nombre, tools permitidas, si es invocable, modelo a usar, agente asociado."),
    ("Cuerpo (Markdown)", "Las instrucciones en lenguaje natural que Claude sigue al invocar la skill."),
]
y = 2.7
for h, d in fields:
    txt(s, 8.0, y, 4.3, 0.4, [[(h, 13.5, CYAN, True)]])
    txt(s, 8.0, y + 0.38, 4.3, 0.9, [[(d, 12, LIGHT, False)]], line_spacing=1.05)
    y += 1.5
txt(s, 8.0, 5.75, 4.3, 0.7,
    [[("1 skill = 1 carpeta con un SKILL.md. Sin build, sin compilar.", 12.5, MUTED, True)]], line_spacing=1.05)

# ================= SLIDE 8 — Profiles =================
s = slide()
chrome(s, 8, "EL SISTEMA LEGO", "Profiles: combinar por stack")
txt(s, 0.78, 1.85, 11.7, 0.5,
    [[("Cada profile es un YAML que agrupa skills, agentes, reglas y permisos. Se combinan con ", 13.5, LIGHT, False),
      ("+", 14, AMBER, True), (".", 13.5, LIGHT, False)]])
# bases
rect(s, 0.78, 2.5, 5.7, 3.0, PANEL, radius=True)
txt(s, 1.0, 2.65, 5.3, 0.4, [[("BASE  (elegis uno)", 13, GREEN, True)]])
bases = [("unity-dev", "25"), ("backend-ts", "33"), ("backend-go", "19"),
         ("frontend", "14"), ("mobile", "12"), ("flutter / ios / android", "8")]
y = 3.15
for name, n in bases:
    txt(s, 1.0, y, 4.2, 0.35, [[(name, 13, LIGHT, False)]])
    txt(s, 5.3, y, 0.9, 0.35, [[(n, 13, CYAN, True)]], align=PP_ALIGN.RIGHT)
    y += 0.37
# addons
rect(s, 6.63, 2.5, 5.85, 3.0, PANEL, radius=True)
txt(s, 6.85, 2.65, 5.4, 0.4, [[("ADDONS  (combinas con +)", 13, AMBER, True)]])
addons = [("+agile", "10"), ("+testing", "11"), ("+infra", "17"),
          ("+security", "5"), ("+marketing", "44"), ("+clevel / +regulatory", "28 / 13")]
y = 3.15
for name, n in addons:
    txt(s, 6.85, y, 4.4, 0.35, [[(name, 13, LIGHT, False)]])
    txt(s, 11.3, y, 0.9, 0.35, [[(n, 13, CYAN, True)]], align=PP_ALIGN.RIGHT)
    y += 0.37
rect(s, 0.78, 5.7, 11.7, 0.85, PANEL2, radius=True)
txt(s, 1.0, 5.78, 11.3, 0.7,
    [[("+ core (siempre): 21 skills universales + 15 hooks + seguridad.  ", 13, WHITE, True),
      ("Profiles combinados nunca duplican un skill.", 13, CYAN, True)]],
    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)

# ================= SLIDE 9 — Mecanismos avanzados =================
s = slide()
chrome(s, 9, "PRODUCTION-GRADE", "Mecanismos que lo elevan")
feats = [
    ("Distribucion hibrida", "Resuelve el contenido en cadena: GitHub -> cache local -> bundled en npm. Funciona incluso sin conexion.", CYAN),
    ("Smart update (3-way merge)", "Compara base / tu version / upstream. Actualiza lo no tocado, preserva tus cambios, avisa conflictos.", VIOLET),
    ("Worktree-aware", "En git worktrees comparte hooks/ y docs/ (symlink), pero copia skills/ por rama (pueden divergir).", GREEN),
    ("Auto-release", "semantic-release lee los conventional commits: fix -> patch, feat -> minor. Publica a npm solo.", AMBER),
]
positions = [(0.78, 1.95), (6.93, 1.95), (0.78, 4.3), (6.93, 4.3)]
cw, ch = 5.6, 2.15
for (h, d, col), (x, y) in zip(feats, positions):
    rect(s, x, y, cw, ch, PANEL, radius=True)
    rect(s, x, y, 0.1, ch, col)
    txt(s, x + 0.32, y + 0.22, cw - 0.6, 0.5, [[(h, 17, col, True)]])
    txt(s, x + 0.32, y + 0.78, cw - 0.62, 1.2, [[(d, 12.5, LIGHT, False)]], line_spacing=1.1)

# ================= SLIDE 10 — Calidad / testing =================
s = slide()
chrome(s, 10, "CONFIABILIDAD", "Las skills se testean")
txt(s, 0.78, 1.9, 11.7, 0.5,
    [[("322 prompts hay que mantenerlos con calidad. Arcane trae un framework de QA self-contained: ", 13.5, LIGHT, False),
      ("skills-selftest/", 13.5, CYAN, True)]])
cards = [
    ("/skill-test static", "Chequeo estructural: frontmatter valido, paths, tools declaradas.", GREEN),
    ("/skill-test spec", "Evaluacion de comportamiento contra casos de prueba esperados.", CYAN),
    ("/skill-test audit", "Reporte de cobertura: que skills tienen specs y cuales faltan.", VIOLET),
    ("/skill-improve", "Loop test -> diagnostica -> arregla -> re-testea -> mantiene o revierte.", AMBER),
]
y = 2.55
for h, d, col in cards:
    rect(s, 0.78, y, 11.7, 0.85, PANEL, radius=True)
    rect(s, 0.78, y, 0.1, 0.85, col)
    txt(s, 1.05, y + 0.08, 3.6, 0.7, [[(h, 15, col, True)]], anchor=MSO_ANCHOR.MIDDLE, font=MONO)
    txt(s, 4.9, y + 0.08, 7.4, 0.7, [[(d, 12.5, LIGHT, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.98
txt(s, 0.78, 6.5, 11.7, 0.4,
    [[("catalog.yaml = registro maestro  ·  quality-rubric.md = criterios pass/fail por categoria.", 12.5, MUTED, True)]])

# ================= SLIDE 11 — Por que gana =================
s = slide()
chrome(s, 11, "DIFERENCIADORES", "Por que Arcane gana")
rect(s, 0.78, 1.95, 5.7, 0.52, PANEL2, radius=True)
rect(s, 6.63, 1.95, 5.85, 0.52, PANEL2, radius=True)
txt(s, 1.0, 1.98, 5.3, 0.45, [[("ALTERNATIVA", 12.5, MUTED, True)]], anchor=MSO_ANCHOR.MIDDLE)
txt(s, 6.85, 1.98, 5.4, 0.45, [[("COMO GANA ARCANE", 12.5, CYAN, True)]], anchor=MSO_ANCHOR.MIDDLE)
rows = [
    ("Config a mano", "Un comando, versionado, reproducible"),
    ("Mega-repo monolitico", "Selectivo por profile + dedup (menos tokens)"),
    ("Copiar config entre proyectos", "Smart update con three-way merge"),
    ("Claude \"vanilla\"", "Hooks + rules + 86 agentes especializados"),
]
y = 2.58
for a, b in rows:
    rect(s, 0.78, y, 5.7, 0.68, PANEL, radius=True)
    rect(s, 6.63, y, 5.85, 0.68, PANEL, radius=True)
    txt(s, 1.0, y + 0.04, 5.3, 0.6, [[(a, 13.5, LIGHT, False)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 6.85, y + 0.04, 5.4, 0.6, [[(b, 13, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.78
rect(s, 0.78, y + 0.05, 11.7, 0.68, VIOLET, radius=True)
txt(s, 0.78, y + 0.08, 11.7, 0.6,
    [[("Configuracion de IA tratada como software: versionada, modular, testeable y distribuida.", 13.5, WHITE, True)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ================= SLIDE 12 — Trade-offs + cierre =================
s = slide()
chrome(s, 12, "BALANCE Y CIERRE", "Trade-offs honestos")
# pros
rect(s, 0.78, 1.95, 5.7, 3.05, PANEL, radius=True)
rect(s, 0.78, 1.95, 5.7, 0.55, RGBColor(0x14, 0x3D, 0x2E), radius=True)
txt(s, 1.02, 1.98, 5.3, 0.5, [[("FORTALEZAS", 14, GREEN, True)]], anchor=MSO_ANCHOR.MIDDLE)
pros = ["Eficiencia de tokens", "Un comando, cero friccion",
        "Reproducible (manifest)", "Estandarizacion de equipo",
        "Auto-update sin pisar cambios", "Testeable (skills-selftest)"]
y = 2.62
for p in pros:
    txt(s, 1.02, y, 5.3, 0.4, [[("+  ", 13, GREEN, True), (p, 12.5, LIGHT, False)]])
    y += 0.39
# contras
rect(s, 6.63, 1.95, 5.85, 3.05, PANEL, radius=True)
rect(s, 6.63, 1.95, 5.85, 0.55, RGBColor(0x4A, 0x39, 0x12), radius=True)
txt(s, 6.87, 1.98, 5.4, 0.5, [[("COSTOS", 14, AMBER, True)]], anchor=MSO_ANCHOR.MIDDLE)
cons = ["Mantener 322 skills con calidad", "Skills = prompts -> no deterministas",
        "Depende de la API de Claude Code", "Curva: saber que profiles combinar"]
y = 2.62
for c in cons:
    txt(s, 6.87, y, 5.5, 0.5, [[("!  ", 13, AMBER, True), (c, 12.5, LIGHT, False)]])
    y += 0.52
# CTA
rect(s, 0.78, 5.25, 11.7, 1.45, VIOLET, radius=True)
txt(s, 0.78, 5.4, 11.7, 0.5,
    [[("Configuracion de IA, modular y versionada — un comando para empezar.", 17, WHITE, True)]],
    align=PP_ALIGN.CENTER)
txt(s, 0.78, 5.95, 11.7, 0.5, [[("npx arcane install backend-ts+agile", 21, WHITE, True)]],
    align=PP_ALIGN.CENTER, font=MONO)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "arcane-overview-12.pptx")
prs.save(out)
print("OK ->", out)
print("Slides:", len(prs.slides._sldIdLst))
